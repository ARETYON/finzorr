"""Sanity: local disk document storage (traversal jail) + text extraction."""

import io
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF
import pytest

from app.core.config import settings
from app.documents.ingest import (
    DocumentTooLargeError,
    UnsupportedDocumentError,
    extract_any,
    extract_pages,
    extract_text_file,
)
from app.documents.storage import DocumentStorage, LocalDiskStorage, get_storage
from app.domain.chunking import CHUNK_CHARS, CHUNK_OVERLAP, chunk_pages

pytestmark = pytest.mark.sanity


# ---------------------------------------------------------------- LocalDiskStorage


async def test_save_load_delete_roundtrip(tmp_path: Path) -> None:
    storage = LocalDiskStorage(str(tmp_path))
    key = "user-1/doc-1.pdf"
    await storage.save(key, b"hello bytes")
    assert await storage.load(key) == b"hello bytes"
    assert (tmp_path / "user-1" / "doc-1.pdf").read_bytes() == b"hello bytes"
    await storage.delete(key)
    assert not (tmp_path / "user-1" / "doc-1.pdf").exists()


async def test_save_overwrites_existing_key(tmp_path: Path) -> None:
    storage = LocalDiskStorage(str(tmp_path))
    await storage.save("a.bin", b"one")
    await storage.save("a.bin", b"two")
    assert await storage.load("a.bin") == b"two"


async def test_delete_missing_key_is_silent(tmp_path: Path) -> None:
    storage = LocalDiskStorage(str(tmp_path))
    await storage.delete("never-saved.pdf")  # missing_ok — no error


async def test_load_missing_key_raises(tmp_path: Path) -> None:
    storage = LocalDiskStorage(str(tmp_path))
    with pytest.raises(FileNotFoundError):
        await storage.load("nope.pdf")


async def test_path_traversal_rejected_on_every_operation(tmp_path: Path) -> None:
    storage = LocalDiskStorage(str(tmp_path / "jail"))
    (tmp_path / "secret.txt").write_bytes(b"outside")
    for key in ("../secret.txt", "a/../../secret.txt"):
        with pytest.raises(ValueError, match="illegal storage key"):
            await storage.save(key, b"x")
        with pytest.raises(ValueError, match="illegal storage key"):
            await storage.load(key)
        with pytest.raises(ValueError, match="illegal storage key"):
            await storage.delete(key)
    assert (tmp_path / "secret.txt").read_bytes() == b"outside"  # untouched


async def test_absolute_key_rejected(tmp_path: Path) -> None:
    storage = LocalDiskStorage(str(tmp_path))
    with pytest.raises(ValueError, match="illegal storage key"):
        await storage.load("/etc/passwd")


def test_get_storage_is_a_local_disk_singleton() -> None:
    first = get_storage()
    assert isinstance(first, LocalDiskStorage)
    assert isinstance(first, DocumentStorage)
    assert get_storage() is first


# ---------------------------------------------------------------- extract_text_file


def test_csv_becomes_readable_rows() -> None:
    pages = extract_text_file(b"name,price\nTCS,100\nINFY,50\n")
    assert len(pages) == 1
    lines = pages[0][1].splitlines()
    assert lines[0] == "Columns: name, price"
    assert lines[1] == "name=TCS; price=100"
    assert lines[2] == "name=INFY; price=50"


def test_single_column_file_stays_plain_text() -> None:
    pages = extract_text_file(b"just\nsome\nlines\n")
    assert pages == [("p.1", "just\nsome\nlines\n")]


def test_plain_text_passthrough_and_empty() -> None:
    assert extract_text_file(b"hello world") == [("p.1", "hello world")]
    assert extract_text_file(b"") == []
    assert extract_text_file(b"   \n\t ") == []


def test_invalid_utf8_is_replaced_not_fatal() -> None:
    pages = extract_text_file(b"caf\xff latte")
    assert len(pages) == 1
    assert "�" in pages[0][1]


# ---------------------------------------------------------------- extract_any routing


def test_extract_any_routes_txt_md_csv() -> None:
    assert extract_any("notes.txt", b"plain") == [("p.1", "plain")]
    assert extract_any("README.md", b"# title") == [("p.1", "# title")]
    label, text = extract_any("Data.CSV", b"a,b\n1,2\n")[0]
    assert label == "p.1" and text.startswith("Columns: a, b")


def test_extract_any_rejects_unknown_types() -> None:
    with pytest.raises(UnsupportedDocumentError, match="virus.exe"):
        extract_any("virus.exe", b"MZ")
    with pytest.raises(UnsupportedDocumentError):
        extract_any("archive.zip", b"PK")


def test_legacy_ppt_gets_actionable_message() -> None:
    with pytest.raises(UnsupportedDocumentError, match="save as .pptx"):
        extract_any("deck.ppt", b"\xd0\xcf\x11\xe0old")


def _pdf_bytes(pages: list[str]) -> bytes:
    doc = fitz.open()
    for text in pages:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    data: bytes = doc.tobytes()
    doc.close()
    return data


def test_extract_any_reads_pdf_pages() -> None:
    extracted = extract_any("report.pdf", _pdf_bytes(["alpha page", "beta page"]))
    assert len(extracted) == 2
    assert extracted[0][0] == "p.1" and "alpha page" in extracted[0][1]
    assert extracted[1][0] == "p.2" and "beta page" in extracted[1][1]


def test_pdf_over_page_cap_rejected(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(settings, "MAX_UPLOAD_PAGES", 1)
    with pytest.raises(DocumentTooLargeError, match="2 pages exceeds the 1-page limit"):
        extract_pages(_pdf_bytes(["one", "two"]))


def _scanned_pdf_bytes(text: str) -> bytes:
    """A page with NO text layer — real content only exists as a rendered
    image, same shape as an actual scanned document (zero-length get_text()).
    """
    src = fitz.open()
    src_page = src.new_page()
    src_page.insert_text((72, 72), text, fontsize=14)
    pix = src_page.get_pixmap(dpi=200)

    scanned = fitz.open()
    img_page = scanned.new_page()
    img_page.insert_image(img_page.rect, pixmap=pix)
    data: bytes = scanned.tobytes()
    scanned.close()
    src.close()
    return data


def test_scanned_pdf_falls_back_to_ocr() -> None:
    pages = extract_pages(_scanned_pdf_bytes("Ruby brings vitality and confidence"))
    assert len(pages) == 1
    label, text = pages[0]
    assert label == "p.1"
    assert "ruby" in text.lower()
    assert "vitality" in text.lower()


def test_normal_pdf_text_layer_skips_ocr(monkeypatch: pytest.MonkeyPatch) -> None:
    """A page with a real, substantial text layer must never trigger OCR --
    OCR is CPU-expensive and the native text layer is already authoritative.
    """

    def _fail_if_called(*_args: Any, **_kwargs: Any) -> str:
        raise AssertionError("OCR should not run when a real text layer exists")

    monkeypatch.setattr("pytesseract.image_to_string", _fail_if_called)
    long_text = "This page has a completely normal, substantial text layer. " * 3
    pages = extract_pages(_pdf_bytes([long_text]))
    assert long_text.strip()[:20] in pages[0][1]


def test_ocr_failure_degrades_gracefully(monkeypatch: pytest.MonkeyPatch) -> None:
    """A missing/broken Tesseract install must not fail the whole upload --
    the page just keeps whatever (possibly empty) native text it had.
    """
    monkeypatch.setattr(
        "pytesseract.image_to_string",
        lambda *_a, **_k: (_ for _ in ()).throw(RuntimeError("tesseract not found")),
    )
    pages = extract_pages(_scanned_pdf_bytes("unreachable without OCR"))
    assert pages == [("p.1", "")]


# ---------------------------------------------------------------- chunk_pages


def test_chunking_windows_with_overlap_and_locators() -> None:
    long_text = "abcdefghij" * 200  # 2000 chars, > CHUNK_CHARS
    chunks = chunk_pages([("p.1", long_text)], "big.txt")
    assert len(chunks) == 2
    assert all(c["title"] == "big.txt" and c["locator"] == "p.1" for c in chunks)
    assert chunks[0]["text"] == long_text[:CHUNK_CHARS]
    step = CHUNK_CHARS - CHUNK_OVERLAP
    assert chunks[1]["text"] == long_text[step : step + CHUNK_CHARS]


def test_blank_pages_skipped_and_whitespace_normalized() -> None:
    chunks = chunk_pages(
        [("p.1", ""), ("p.2", "   \n\t "), ("p.3", "hello   world\n\nagain")], "doc.pdf"
    )
    assert len(chunks) == 1
    assert chunks[0]["text"] == "hello world again"
    assert chunks[0]["locator"] == "p.3"  # labels survive blank-page skipping


def test_short_page_is_single_chunk() -> None:
    chunks = chunk_pages([("p.1", "tiny")], "t.txt")
    assert chunks == [{"text": "tiny", "title": "t.txt", "locator": "p.1"}]


# ------------------------------------------------- new formats: pptx / xlsx / xls


def _pptx_bytes(slides: list[str]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    blank = deck.slide_layouts[6]
    for text in slides:
        slide = deck.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = text
    buf = io.BytesIO()
    deck.save(buf)
    return buf.getvalue()


def _xlsx_bytes(sheets: dict[str, list[list[str]]]) -> bytes:
    import openpyxl

    book = openpyxl.Workbook()
    book.remove(book.active)
    for name, rows in sheets.items():
        sheet = book.create_sheet(title=name)
        for row in rows:
            sheet.append(row)
    buf = io.BytesIO()
    book.save(buf)
    return buf.getvalue()


def test_pptx_slides_become_labeled_pages() -> None:
    pages = extract_any("deck.pptx", _pptx_bytes(["intro slide", "revenue grew 20%"]))
    assert [label for label, _ in pages] == ["slide 1", "slide 2"]
    assert "revenue grew 20%" in pages[1][1]


def test_pptx_over_slide_cap_rejected(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(settings, "MAX_UPLOAD_PAGES", 1)
    with pytest.raises(DocumentTooLargeError, match="2 slides exceeds"):
        extract_any("deck.pptx", _pptx_bytes(["one", "two"]))


def test_xlsx_sheets_become_labeled_pages() -> None:
    pages = extract_any(
        "book.xlsx",
        _xlsx_bytes(
            {
                "Revenue": [["quarter", "amount"], ["Q3", "500"]],
                "Costs": [["item", "value"], ["rent", "90"]],
            }
        ),
    )
    assert [label for label, _ in pages] == ["sheet:Revenue", "sheet:Costs"]
    assert "quarter=Q3; amount=500" in pages[0][1]
    assert pages[0][1].startswith("Columns: quarter, amount")


def test_docx_sections_labeled() -> None:
    import docx as docx_module

    document = docx_module.Document()
    document.add_paragraph("hello world")
    buf = io.BytesIO()
    document.save(buf)
    pages = extract_any("note.docx", buf.getvalue())
    assert pages[0][0] == "§1"
    assert "hello world" in pages[0][1]


# ------------------------------------------------ orphan cleanup on failure


async def test_failed_ingest_cleans_up_attempted_batches(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Embed dies on batch 2 -> the batch-0 vectors already upserted (and the
    attempted batch-1 key) are best-effort deleted before the error
    propagates — otherwise they orphan forever (failed docs have
    chunk_count=NULL and re-uploads mint a new doc id)."""
    import uuid as uuid_module

    from app.documents import ingest as ingest_mod

    calls = {"embed": 0}
    deleted: list[tuple[str, str]] = []

    async def flaky_embed(texts: list[str]) -> list[list[float]]:
        calls["embed"] += 1
        if calls["embed"] >= 2:
            raise RuntimeError("embedder down")
        return [[0.0] * 3 for _ in texts]

    async def fake_upsert(tenant: str, doc_id: str, batch: list[Any], vectors: list[Any]) -> int:
        return len(batch)

    async def record_delete(tenant: str, doc_id: str) -> None:
        deleted.append((tenant, doc_id))

    monkeypatch.setattr(ingest_mod, "embed_texts", flaky_embed)
    monkeypatch.setattr(ingest_mod, "upsert_chunks", fake_upsert)
    monkeypatch.setattr("app.infrastructure.vector_store.delete_document", record_delete)

    user_id, doc_id = uuid_module.uuid4(), uuid_module.uuid4()
    # >16 chunks => at least 2 embed batches (1200-char chunks from 1 page)
    big_page = "word " * 8000
    with pytest.raises(RuntimeError, match="embedder down"):
        await ingest_mod.ingest_document(user_id, doc_id, "big.txt", big_page.encode())

    assert deleted, "attempted batches must be cleaned up"
    assert all(t == str(user_id) for t, _ in deleted)
    assert f"{doc_id}:0" in [d for _, d in deleted]  # the upserted batch
    assert f"{doc_id}:16" in [d for _, d in deleted]  # the attempted batch
