"""Document -> chunks -> embeddings -> Qdrant (per-user tenant).

Supported: PDF (PyMuPDF), DOCX (python-docx), PPTX (python-pptx),
XLSX (openpyxl), legacy XLS (xlrd), CSV/TXT/MD (plain decode). Legacy .ppt
has no viable free parser — rejected with a "save as .pptx" message. OCR
for scanned PDFs is a Phase-2 addition. Extractors return LABELED
pseudo-pages [(label, text)] so citations read naturally per format:
`file.pdf · p.3`, `deck.pptx · slide 4`, `book.xlsx · sheet:Revenue`.
"""

import asyncio
import csv as csv_module
import io
import uuid

import fitz  # PyMuPDF
from langsmith import traceable

from app.core.config import settings
from app.core.logging import log
from app.core.trace import mark
from app.rag.embeddings import embed_texts
from app.rag.vector_store import upsert_chunks

CHUNK_CHARS = 1200
CHUNK_OVERLAP = 200
_EMBED_BATCH = 16
_CSV_MAX_ROWS = 2000


class DocumentTooLargeError(Exception):
    """Raised when a document exceeds the configured page/size caps."""


class UnsupportedDocumentError(Exception):
    """Raised when a file type cannot be extracted."""


LabeledPages = list[tuple[str, str]]  # (locator label, text)


def extract_pages(pdf_bytes: bytes) -> LabeledPages:
    """Extract text per page; raises DocumentTooLargeError over the page cap."""
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        if doc.page_count > settings.MAX_UPLOAD_PAGES:
            raise DocumentTooLargeError(
                f"{doc.page_count} pages exceeds the {settings.MAX_UPLOAD_PAGES}-page limit"
            )
        return [(f"p.{i}", page.get_text()) for i, page in enumerate(doc, start=1)]


def chunk_pages(pages: LabeledPages, filename: str) -> list[dict[str, str]]:
    """Window each labeled pseudo-page, keeping its label as the citation
    anchor (`p.N` / `slide N` / `sheet:Name` / `§N`)."""
    chunks: list[dict[str, str]] = []
    for label, text in pages:
        clean = " ".join(text.split())
        if not clean:
            continue
        start = 0
        while start < len(clean):
            piece = clean[start : start + CHUNK_CHARS]
            chunks.append({"text": piece, "title": filename, "locator": label})
            if start + CHUNK_CHARS >= len(clean):
                break
            start += CHUNK_CHARS - CHUNK_OVERLAP
    return chunks


def extract_docx(data: bytes) -> LabeledPages:
    """Extract DOCX text as sections (~40 paragraphs each), cited as §N."""
    import docx

    document = docx.Document(io.BytesIO(data))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            paragraphs.append(" | ".join(cell.text.strip() for cell in row.cells))
    if not paragraphs:
        return []
    section_size = 40
    return [
        (f"§{n}", "\n".join(paragraphs[i : i + section_size]))
        for n, i in enumerate(range(0, len(paragraphs), section_size), start=1)
    ]


def extract_pptx(data: bytes) -> LabeledPages:
    """One pseudo-page per slide (shape text + tables), cited as `slide N`."""
    from pptx import Presentation

    deck = Presentation(io.BytesIO(data))
    slides = list(deck.slides)
    if len(slides) > settings.MAX_UPLOAD_PAGES:
        raise DocumentTooLargeError(
            f"{len(slides)} slides exceeds the {settings.MAX_UPLOAD_PAGES}-slide limit"
        )
    pages: LabeledPages = []
    for n, slide in enumerate(slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.extend(
                    para.text for para in shape.text_frame.paragraphs if para.text.strip()
                )
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        pages.append((f"slide {n}", "\n".join(parts)))
    return pages


_SHEET_MAX_ROWS = 2000


def _rows_to_text(header: list[str], rows: list[list[str]]) -> str:
    """Same readable flattening the CSV path uses: `col=value; col=value`."""
    lines = [f"Columns: {', '.join(header)}"] if header else []
    lines += [
        "; ".join(f"{h}={v}" for h, v in zip(header, row, strict=False))
        for row in rows[:_SHEET_MAX_ROWS]
    ]
    return "\n".join(lines)


def extract_xlsx(data: bytes) -> LabeledPages:
    """One pseudo-page per sheet, cited as `sheet:Name` (row-capped)."""
    import openpyxl

    book = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        if len(book.sheetnames) > settings.MAX_UPLOAD_PAGES:
            raise DocumentTooLargeError(
                f"{len(book.sheetnames)} sheets exceeds the "
                f"{settings.MAX_UPLOAD_PAGES}-sheet limit"
            )
        pages: LabeledPages = []
        for sheet in book.worksheets:
            rows = [
                ["" if c is None else str(c) for c in row]
                for row in sheet.iter_rows(values_only=True, max_row=_SHEET_MAX_ROWS + 1)
            ]
            header, body = (rows[0], rows[1:]) if rows else ([], [])
            pages.append((f"sheet:{sheet.title}", _rows_to_text(header, body)))
        return pages
    finally:
        book.close()


def extract_xls(data: bytes) -> LabeledPages:
    """Legacy .xls via xlrd — same sheet-per-page shape as XLSX."""
    import xlrd

    book = xlrd.open_workbook(file_contents=data)
    if book.nsheets > settings.MAX_UPLOAD_PAGES:
        raise DocumentTooLargeError(
            f"{book.nsheets} sheets exceeds the {settings.MAX_UPLOAD_PAGES}-sheet limit"
        )
    pages: LabeledPages = []
    for sheet in book.sheets():
        rows = [
            [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
            for r in range(min(sheet.nrows, _SHEET_MAX_ROWS + 1))
        ]
        header, body = (rows[0], rows[1:]) if rows else ([], [])
        pages.append((f"sheet:{sheet.name}", _rows_to_text(header, body)))
    return pages


def extract_text_file(data: bytes) -> LabeledPages:
    """Plain text/CSV decode as one pseudo-page (CSV rows joined readably)."""
    text = data.decode("utf-8", errors="replace")
    try:
        rows = list(csv_module.reader(io.StringIO(text)))
        if len(rows) > 1 and len(rows[0]) > 1:  # looks like a real CSV
            header = ", ".join(rows[0])
            lines = [f"Columns: {header}"] + [
                "; ".join(f"{h}={v}" for h, v in zip(rows[0], r, strict=False))
                for r in rows[1 : _CSV_MAX_ROWS + 1]
            ]
            text = "\n".join(lines)
    except csv_module.Error:
        pass
    return [("p.1", text)] if text.strip() else []


def extract_any(filename: str, data: bytes) -> LabeledPages:
    """Route extraction by extension; raises UnsupportedDocumentError otherwise."""
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return extract_pages(data)
    if lower.endswith(".docx"):
        return extract_docx(data)
    if lower.endswith(".pptx"):
        return extract_pptx(data)
    if lower.endswith(".xlsx"):
        return extract_xlsx(data)
    if lower.endswith(".xls"):
        return extract_xls(data)
    if lower.endswith(".ppt"):
        raise UnsupportedDocumentError(
            "legacy .ppt isn't supported — please save as .pptx and re-upload"
        )
    if lower.endswith((".csv", ".txt", ".md")):
        return extract_text_file(data)
    raise UnsupportedDocumentError(f"unsupported file type: {filename}")


_LOADERS = {
    ".pdf": "pymupdf",
    ".docx": "python-docx",
    ".pptx": "python-pptx",
    ".xlsx": "openpyxl",
    ".xls": "xlrd",
}


def _loader_for(filename: str) -> str:
    lower = filename.lower()
    for ext, loader in _LOADERS.items():
        if lower.endswith(ext):
            return loader
    return "plain-text"


@traceable(
    run_type="parser",
    name="extract",
    process_outputs=lambda pages: {
        "pages": len(pages),
        "labels": [label for label, _ in pages][:20],
        "total_chars": sum(len(t) for _, t in pages),
    },
)
def _extract_traced(filename: str, data: bytes) -> LabeledPages:
    mark(loader=_loader_for(filename))
    return extract_any(filename, data)


@traceable(
    run_type="parser",
    name="split",
    process_inputs=lambda inputs: {"pages": len(inputs.get("pages", [])),
                                   "filename": inputs.get("filename", "")},
    process_outputs=lambda chunks: {
        "chunks": len(chunks),
        "chunk_chars": CHUNK_CHARS,
        "overlap": CHUNK_OVERLAP,
    },
)
def _split_traced(pages: LabeledPages, filename: str) -> list[dict[str, str]]:
    return chunk_pages(pages, filename)


@traceable(run_type="chain", name="document.ingest", tags=["rag", "ingest"])
async def ingest_document(
    user_id: uuid.UUID, doc_id: uuid.UUID, filename: str, data: bytes
) -> int:
    """Extract (any supported type), chunk, embed, upsert. Returns chunk count.

    A mid-loop failure (embed outage, upsert error) must not strand the
    batches already upserted: the router marks the doc `failed` with
    chunk_count=NULL, its delete path then only clears batch 0, and a
    re-upload mints a NEW doc id — so ingest-time cleanup is the only
    chance to remove them.
    """
    mark(filename=filename, format=filename.rsplit(".", 1)[-1].lower(), bytes=len(data))
    pages = await asyncio.to_thread(_extract_traced, filename, data)
    chunks = _split_traced(pages, filename)
    if not chunks:
        return 0
    total = 0
    attempted_batches: list[str] = []
    try:
        for start in range(0, len(chunks), _EMBED_BATCH):
            batch = chunks[start : start + _EMBED_BATCH]
            # record BEFORE the attempt — a failed upsert's server-side
            # state must be assumed written, not guessed clean
            attempted_batches.append(f"{doc_id}:{start}")
            vectors = await embed_texts([c["text"] for c in batch])
            total += await upsert_chunks(str(user_id), f"{doc_id}:{start}", batch, vectors)
    except Exception:
        from app.rag.vector_store import delete_document as qdrant_delete

        for batch_key in attempted_batches:
            try:
                await qdrant_delete(str(user_id), batch_key)
            except Exception as cleanup_exc:  # noqa: BLE001 — best-effort
                log.warning(
                    "document.orphan_cleanup_failed",
                    doc_id=str(doc_id),
                    batch=batch_key,
                    error=str(cleanup_exc),
                )
        raise
    log.info("document.ingested", doc_id=str(doc_id), chunks=total, pages=len(pages))
    return total
